"""Generate the generic, fact-bounded Word and PowerPoint handoff files.

The repository does not contain the school's thesis or defense templates, so the
outputs are deliberately labelled as drafts awaiting the school's template.
"""

from __future__ import annotations

import re
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
THESIS = ROOT / "docs" / "thesis"
OUT = THESIS / "deliverables"


def strip_markdown(text: str) -> str:
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
    text = text.replace("**", "").replace("`", "")
    return text.strip()


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_run_font(run, name: str = "宋体", size: float = 10.5, bold: bool = False) -> None:
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    run.font.size = Pt(size)
    run.bold = bold


def add_doc_paragraph(doc: Document, text: str, style: str | None = None, bold: bool = False) -> None:
    para = doc.add_paragraph(style=style)
    run = para.add_run(strip_markdown(text))
    set_run_font(run, size=10.5, bold=bold)
    para.paragraph_format.line_spacing = 1.35
    para.paragraph_format.space_after = Pt(5)


def add_word_deliverable() -> Path:
    draft_path = THESIS / "RIKE_THESIS_DRAFT.md"
    fact_path = THESIS / "RIKE_THESIS_FACT_CHECK.md"
    refs_path = ROOT / "docs" / "THESIS_REFERENCES.md"

    doc = Document()
    section = doc.sections[0]
    section.top_margin = Cm(2.54)
    section.bottom_margin = Cm(2.54)
    section.left_margin = Cm(3.0)
    section.right_margin = Cm(2.5)

    normal = doc.styles["Normal"]
    normal.font.name = "宋体"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    normal.font.size = Pt(10.5)

    for style_name, size, color in (("Title", 22, "17365D"), ("Heading 1", 16, "17365D"), ("Heading 2", 13, "1F4E79"), ("Heading 3", 11, "1F4E79")):
        style = doc.styles[style_name]
        style.font.name = "黑体"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "黑体")
        style.font.size = Pt(size)
        style.font.color.rgb = RGBColor.from_string(color)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.paragraph_format.space_before = Pt(95)
    run = title.add_run("面向高中物化生的 Spring Boot 大模型题库系统设计与实现")
    set_run_font(run, "黑体", 22, True)
    run.font.color.rgb = RGBColor(23, 54, 93)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.paragraph_format.space_before = Pt(18)
    run = subtitle.add_run("RIKE 理科学习辅助系统\n论文事实稿与交付摘要")
    set_run_font(run, "宋体", 15, False)

    notice = doc.add_paragraph()
    notice.alignment = WD_ALIGN_PARAGRAPH.CENTER
    notice.paragraph_format.space_before = Pt(80)
    run = notice.add_run("通用排版草稿：待套用学校论文 Word 模板\n版本日期：2026-08-16")
    set_run_font(run, "宋体", 11, False)
    run.font.color.rgb = RGBColor(96, 96, 96)

    doc.add_page_break()
    add_doc_paragraph(doc, "交付边界", "Heading 1")
    add_doc_paragraph(doc, "本文件由当前仓库事实稿生成，保留实现、数据库、测试和 Provider 边界，不替代学校封面、页眉页脚、目录、格式规范或导师要求。仓库中未发现学校论文 Word 模板，因此本文件必须在提交前套用学校模板并进行人工排版复核。")
    add_doc_paragraph(doc, "事实截止于 2026-08-16。PR #33 仍为 Draft、OPEN、未合并；用户真人最终复验仍待完成。机器浏览器、Mock/Fake Provider 与自动化测试不等同于真人验收或真实 Provider PASS。")

    add_doc_paragraph(doc, "目录（套用学校模板后请在 Word 中更新目录域）", "Heading 1")
    for label in ["摘要与关键词", "第一章 绪论", "第二章 需求分析", "第三章 系统总体设计", "第四章 数据库设计", "第五章 核心功能实现", "第六章 AI Provider 与受控生成", "第七章 权限与安全", "第八章 测试与结果", "第九章 总结与展望", "参考文献", "事实核验附录"]:
        add_doc_paragraph(doc, label)

    doc.add_page_break()
    lines = draft_path.read_text(encoding="utf-8").splitlines()
    in_meta = True
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith(">"):
            continue
        if stripped.startswith("# "):
            in_meta = False
            continue
        if stripped.startswith("## "):
            add_doc_paragraph(doc, stripped[3:], "Heading 1")
            continue
        if stripped.startswith("### "):
            add_doc_paragraph(doc, stripped[4:], "Heading 2")
            continue
        if stripped.startswith("- "):
            p = doc.add_paragraph(style="List Bullet")
            run = p.add_run(strip_markdown(stripped[2:]))
            set_run_font(run)
            p.paragraph_format.line_spacing = 1.25
            continue
        if re.match(r"^\|", stripped):
            # Markdown tables are retained as readable text; the fact check below
            # provides the audit-friendly tabular version.
            add_doc_paragraph(doc, stripped.replace("|", "  "))
            continue
        add_doc_paragraph(doc, stripped)

    doc.add_page_break()
    add_doc_paragraph(doc, "参考文献（正式白名单 22 条）", "Heading 1")
    refs = []
    for line in refs_path.read_text(encoding="utf-8").splitlines():
        match = re.match(r"^\[(\d+)\]\s+(.+)$", line.strip())
        if match:
            refs.append((int(match.group(1)), strip_markdown(match.group(2))))
    if len(refs) != 22:
        raise RuntimeError(f"expected 22 formal references, got {len(refs)}")
    for number, reference in refs:
        add_doc_paragraph(doc, f"[{number}] {reference}")

    doc.add_page_break()
    add_doc_paragraph(doc, "事实核验附录", "Heading 1")
    for line in fact_path.read_text(encoding="utf-8").splitlines():
        stripped = line.strip()
        if not stripped or stripped.startswith(">") or stripped.startswith("# "):
            continue
        if stripped.startswith("## "):
            add_doc_paragraph(doc, stripped[3:], "Heading 2")
        else:
            add_doc_paragraph(doc, stripped.replace("|", "  "))

    props = doc.core_properties
    props.title = "RIKE 论文事实稿（待套学校模板）"
    props.subject = "面向高中物化生的 Spring Boot 大模型题库系统设计与实现"
    props.author = "RIKE 项目交付资料"
    output = OUT / "RIKE_论文事实稿_待套学校模板.docx"
    doc.save(output)
    return output


PPT_NAVY = PptRGBColor(23, 54, 93)
PPT_BLUE = PptRGBColor(31, 78, 121)
PPT_TEAL = PptRGBColor(0, 128, 128)
PPT_BG = PptRGBColor(245, 248, 251)
PPT_TEXT = PptRGBColor(35, 43, 50)
PPT_MUTED = PptRGBColor(95, 105, 115)


def add_box(slide, x, y, w, h, text, *, size=18, color=PPT_TEXT, bold=False, fill=None, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = PptInches(0.12)
    tf.margin_right = PptInches(0.12)
    tf.margin_top = PptInches(0.08)
    tf.margin_bottom = PptInches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = PptPt(2)
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_title(slide, title, source="事实截止：2026-08-16；通用答辩草稿，待套学校模板"):
    add_box(slide, 0.55, 0.28, 12.25, 0.52, title, size=25, color=PPT_NAVY, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.55), PptInches(0.92), PptInches(12.25), PptInches(0.035))
    line.fill.solid(); line.fill.fore_color.rgb = PPT_TEAL; line.line.fill.background()
    add_box(slide, 0.55, 7.08, 12.25, 0.22, source, size=8.5, color=PPT_MUTED)


def add_bullets(slide, items, x=0.8, y=1.35, w=11.6, h=5.25, size=20):
    shape = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = shape.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = PptInches(0.08); tf.margin_right = PptInches(0.08)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = PptPt(size)
        p.font.color.rgb = PPT_TEXT
        p.space_after = PptPt(12)
        p.line_spacing = 1.12
    return shape


def add_image_card(slide, image_path: Path, x, y, w, h, caption):
    slide.shapes.add_picture(str(image_path), PptInches(x), PptInches(y), width=PptInches(w), height=PptInches(h))
    add_box(slide, x, y + h + 0.05, w, 0.35, caption, size=11, color=PPT_MUTED, align=PP_ALIGN.CENTER)


def add_ppt_deliverable() -> Path:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide(title, source=None):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = PPT_BG
        add_title(slide, title, source or "事实截止：2026-08-16；通用答辩草稿，待套学校模板")
        return slide

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = PPT_NAVY
    add_box(slide, 0.8, 1.35, 11.8, 1.0, "面向高中物化生的 Spring Boot 大模型题库系统设计与实现", size=30, color=PptRGBColor(255,255,255), bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 1.2, 2.65, 11.0, 0.6, "RIKE 理科学习辅助系统", size=25, color=PptRGBColor(179, 235, 235), bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 2.0, 4.3, 9.3, 0.9, "本科毕业设计答辩资料\n通用草稿：待套用学校答辩 PPT 模板", size=19, color=PptRGBColor(240,244,248), align=PP_ALIGN.CENTER)
    add_box(slide, 0.8, 6.75, 11.8, 0.25, "版本日期：2026-08-16｜PR #33 Draft / OPEN / 未合并｜用户最终真人复验待完成", size=10, color=PptRGBColor(210,220,230), align=PP_ALIGN.CENTER)

    slide = new_slide("一、研究背景与问题")
    add_bullets(slide, [
        "高中物化生题目需要结构化管理、确定性判分和可追溯的学习反馈。",
        "传统题库难以同时覆盖错题复习、专题学习、知识卡片和教师审核。",
        "大模型可辅助解释与变式生成，但不能替代正式答案、权限边界和人工审核。",
        "研究目标：在前后端分离、模块化单体架构中建立可复验的受控学习闭环。",
    ])

    slide = new_slide("二、研究范围与交付边界")
    add_bullets(slide, [
        "已实现：认证与角色、题库导入审核、练习与错题、专题单元、试卷、知识卡片。",
        "已实现：DeepSeek 文本、GLM/xAI 视觉、Web Search 的统一 Provider 契约与安全降级。",
        "研究证据：代码、Flyway V1–V29、50 张业务表、随机临时库、匿名 Demo、正式库只做受控内容核验。",
        "不声称：真实学校对照实验、学习成绩提升、问卷满意度、AI 学科准确率或本轮 Provider PASS。",
    ])

    slide = new_slide("三、总体架构：前后端分离 + 模块化单体")
    boxes = [(0.8, "Vue 3\nTypeScript\nElement Plus\nPinia / Router"), (4.55, "Spring Boot 4.1\nSecurity + JWT\nMyBatis-Plus\n业务模块服务"), (8.3, "MySQL 8.4\nFlyway V1–V29\n50 张业务表\n正式库 / Demo 隔离")]
    for x, text in boxes:
        add_box(slide, x, 2.0, 3.05, 2.3, text, size=20, color=PPT_NAVY, bold=True, fill=PptRGBColor(225, 237, 245), align=PP_ALIGN.CENTER)
    add_box(slide, 3.55, 2.65, 0.75, 0.55, "→", size=32, color=PPT_TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 7.3, 2.65, 0.75, 0.55, "→", size=32, color=PPT_TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 1.0, 5.2, 11.2, 0.8, "外部 Provider 通过统一 Core 接入；业务结果由状态机、权限检查、确定性判分和安全渲染约束。", size=21, color=PPT_TEXT, fill=PptRGBColor(232, 244, 241), align=PP_ALIGN.CENTER)

    slide = new_slide("四、学生学习闭环：从练习到再学习")
    flow = ["题库 / 专题单元", "练习快照", "确定性判分", "错题筛选与再做", "掌握度与推荐"]
    for i, label in enumerate(flow):
        x = 0.55 + i * 2.55
        add_box(slide, x, 2.2, 2.0, 1.25, label, size=18, color=PPT_NAVY, bold=True, fill=PptRGBColor(225, 237, 245), align=PP_ALIGN.CENTER)
        if i < len(flow)-1:
            add_box(slide, x + 2.0, 2.52, 0.5, 0.5, "→", size=24, color=PPT_TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 1.0, 4.65, 11.1, 1.0, "错题列表不直接移出；答错重做保持活跃，答对后由学生确认是否归档，避免把一次操作误写成掌握。", size=21, color=PPT_TEXT, fill=PptRGBColor(255, 244, 220), align=PP_ALIGN.CENTER)

    slide = new_slide("五、AI 受控链：辅助生成，不覆盖业务事实")
    flow = ["事实输入\n题目 / 答题 / 知识点", "Prompt + Provider\n超时与安全分类", "Schema V2\n字段级 Parser", "新颖度 + 判分\n事务写入", "学生 DRAFT\n提交后 PENDING"]
    for i, label in enumerate(flow):
        x = 0.4 + i * 2.55
        add_box(slide, x, 1.85, 2.1, 1.55, label, size=17, color=PPT_NAVY, bold=True, fill=PptRGBColor(232, 244, 241) if i == 4 else PptRGBColor(225, 237, 245), align=PP_ALIGN.CENTER)
        if i < len(flow)-1:
            add_box(slide, x + 2.08, 2.32, 0.45, 0.5, "→", size=24, color=PPT_TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 1.0, 4.55, 11.1, 1.0, "STANDARD 与正式判分保持确定性；reasoning_content 不进入业务 DTO、消息表或日志；教师/管理员审核是发布前最后一道边界。", size=20, color=PPT_TEXT, fill=PptRGBColor(255, 244, 220), align=PP_ALIGN.CENTER)

    slide = new_slide("六、正式数据与权限边界")
    add_bullets(slide, [
        "正式数据库：rike_tiku；Flyway V29；50 张业务表；本轮不新增迁移。",
        "正式内容核验：6 个已发布专题单元、18 条单元题目关系；每单元 3 道题。",
        "高频考点：65 张已发布卡片，其中 60 张来自结构化内容源，物理/化学/生物各 10 个条目。",
        "学生候选：本人可见 DRAFT；显式提交后才进入教师审核 PENDING；教师查询不到其他学生 DRAFT。",
    ])

    formal_dir = ROOT / "docs" / "evidence" / "pr33-formal-student"
    slide = new_slide("七、正式学生端机器证据：错题与专题")
    for img, x, caption in [("student-wrong-questions.png", 0.65, "错题筛选：无日期筛选、无列表直接移出"), ("student-topic-units.png", 6.75, "专题单元：单元为一级入口")]:
        path = formal_dir / img
        if path.exists():
            add_image_card(slide, path, x, 1.4, 5.8, 4.75, caption)
    add_box(slide, 0.85, 6.45, 11.6, 0.35, "正式浏览器：4 条路线、19 项断言、0 console/page/request error、0 横向溢出；不等同真人验收。", size=13, color=PPT_TEXT, align=PP_ALIGN.CENTER)

    slide = new_slide("八、正式学生端机器证据：高频考点")
    image = formal_dir / "student-high-frequency-points.png"
    if image.exists():
        add_image_card(slide, image, 0.7, 1.25, 7.1, 5.25, "物化生高频考点与二级结论")
    add_bullets(slide, [
        "入口从“知识卡片”调整为“高频考点”。",
        "展示科学内容、公式/条件、推导、例子、易错点和记忆提示。",
        "不伪造考试年份、频次统计或“AI 自动生成练习”承诺。",
    ], x=8.05, y=1.65, w=4.6, h=4.3, size=17)

    slide = new_slide("九、自动化回归与可复验环境")
    add_bullets(slide, [
        "后端专项：39 tests，0 failures，0 errors；学生前端专项：9 files / 30 tests。",
        "历史全量基线：后端 210 tests、前端 214 tests；最终集中全量结果以本轮回归记录为准。",
        "前端 type-check/build 通过，build 保留已知大 chunk warning；正式数据库不作为集成测试数据源。",
        "隔离端口：前端 18080、后端 18081；浏览器使用独立临时 profile，结束后释放端口。",
    ])

    slide = new_slide("十、Provider、安全与研究边界")
    add_bullets(slide, [
        "本轮没有可安全使用的轮换后凭据；DeepSeek variant、DeepSeek tutor、GLM Vision、xAI Vision、Web Search 均为 BLOCKED_EXTERNAL_PROVIDER。",
        "未读取、输出或复述 Key；未把 Mock/Fake 或历史窗口写成真实 PASS。",
        "JWT、密码、Authorization、Base64、Prompt、reasoning_content 和绝对路径不进入论文截图。",
        "论文只报告工程验证事实，不报告未经执行的真实用户效果或 Provider 学科准确率。",
    ])

    slide = new_slide("十一、工作总结与后续")
    add_bullets(slide, [
        "贡献：形成覆盖题库、练习、错题、专题、卡片、试卷和受控 AI 的完整工程闭环。",
        "方法：用 Flyway、权限边界、确定性判分、Parser、事务和匿名证据控制系统风险。",
        "局限：缺少真实学校用户的长期对照数据；Provider 受凭据、额度和网络条件影响。",
        "当前状态：等待用户完成最终真人验收，并由 GPT 独立审查；未获确认前不允许 ordinary merge。",
    ])

    prs.core_properties.title = "RIKE 答辩 PPT（待套学校模板）"
    prs.core_properties.subject = "面向高中物化生的 Spring Boot 大模型题库系统设计与实现"
    output = OUT / "RIKE_答辩PPT_待套学校模板.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    OUT.mkdir(parents=True, exist_ok=True)
    word = add_word_deliverable()
    ppt = add_ppt_deliverable()
    print(f"WORD_OUTPUT={word}")
    print(f"PPT_OUTPUT={ppt}")
